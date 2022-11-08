from docx import Document
from pathlib import Path
from pptx import Presentation
from os.path import join
from werkzeug.utils import secure_filename
from PyPDF2 import PdfFileReader
from nltk.corpus import stopwords
import re

def use_regex(input_text):
    text = re.sub("[^\x00-\x7F]+"," ", input_text)
    textList = re.split("\W+", text)
    return textList

def use_regex_test(input_text):
    text = re.sub("\s+", ",", input_text)
    text = re.sub("[^\x00-\x7F]+", "", text)
    text = text.split(",")
    for i in text:
        if len(text[i]) == 0:
            continue
    return text

def PdfToText (extLocation:str, fileobj: object):
    FileName = Path(fileobj.filename).stem + '.txt'
    pdf = PdfFileReader(fileobj)
    with Path(join(extLocation,secure_filename(FileName))).open(mode="w", encoding='utf-8') as outputFile:
        text = ''
        for page in pdf.pages:
            text += page.extract_text()
        text = use_regex(text)
        return outputFile.write(text)

def DocxToText(extLocation:str, fileobj: object):
    FileName = Path(fileobj.filename).stem + '.txt'
    document = Document(fileobj)
    with Path(join(extLocation,secure_filename(FileName))).open(mode="w", encoding='utf-8') as outputFile:
        docText = '\n\n'.join(paragraph.text + "," for paragraph in document.paragraphs)
        # docText = use_regex(docText)
        return outputFile.write(docText)

def pptxToTxt(fileobj: str):
    prs = Presentation(fileobj)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                # print(shape.text)
                for run in paragraph.runs:
                    text += run.text + " "
    text = use_regex(text)
    text = text.split(",")
    return text
text = 'Basic,Text,ProcessingRegular,ExpressionsRegular,expressionsA,formal,language,for,specifying,text,stringsHow,can,we,search,for,any,of,these,woodchuckwoodchucksWoodchuckWoodchucks,Regular,Expressions,DisjunctionsLetters,inside,square,brackets,Ranges,A,Z,PatternMatches,wW,oodchuckWoodchuck,woodchuck,1234567890,Any,digitPatternMatches,A,Z,An,upper,case,letterDrenched,Blossoms,a,z,A,lower,case,lettermy,beans,were,impatient,0,9,A,singledigitChapter,1,Down,the,Rabbit,HoleRegular,Expressions,Negation,in,DisjunctionNegations,Ss,Carat,means,negation,only,when,first,in,PatternMatches,A,Z,Notan,upper,case,letterOyfnpripetchik,Ss,Neither,S,nor,sIhave,no,exquisite,reason,e,Neither,e,nor,Look,herea,bThe,patternacaratbLook,up,a,bnowRegular,Expressions,More,DisjunctionWoodchucks,is,another,name,for,groundhog,The,pipe,for,disjunctionPatternMatchesgroundhog,woodchuckyours,mineyoursminea,b,c,abc,gG,roundhog,Ww,oodchuck,Regular,Expressions,Stephen,C,KleenePatternMatchescolou,rOptionalprevious,charcolorcolouroo,h,0,or,more,ofprevious,charoh,ooh,oooh,ooooh,o,h,1,or,more,of,previous,charoh,ooh,oooh,ooooh,baa,baabaaabaaaabaaaaabeg,nbegin,begun,begun,beg3nKleene,Kleene,Regular,Expressions,Anchors,PatternMatches,A,Z,PaloAlto,A,Za,z,1Hello,The,end,The,end,The,end,ExampleFind,me,all,instances,of,the,word,the,in,a,text,theMisses,capitalized,examples,tT,heIncorrectly,returns,otheror,theology,a,zA,Z,tT,he,a,zA,Z,ErrorsThe,process,we,just,went,through,was,based,on,fixing,two,kinds,of,errorsMatching,strings,that,we,should,not,have,matched,there,then,other,False,positives,Type,I,Not,matching,things,that,we,should,have,matched,The,False,negatives,Type,II,Errors,cont,In,NLP,we,are,always,dealing,with,these,kinds,of,errors,Reducing,the,error,rate,for,an,application,often,involves,two,antagonistic,efforts,Increasing,accuracy,or,precision,minimizing,false,positives,Increasing,coverage,or,recall,minimizing,false,negatives,SummaryRegular,expressions,play,a,surprisingly,large,roleSophisticated,sequences,of,regular,expressions,are,often,the,first,model,for,any,text,processing,textFor,many,hard,tasks,we,use,machine,learning,classifiersBut,regular,expressions,are,used,as,features,in,the,classifiersCan,be,very,useful,in,capturing,generalizations11Basic,Text,ProcessingRegular,ExpressionsBasic,Text,ProcessingWord,tokenizationText,NormalizationEvery,NLP,task,needs,to,do,text,normalization,1,Segmenting,tokenizing,words,in,running,text2,Normalizing,word,formats3,Segmenting,sentences,in,running,textHow,many,words,I,do,uh,main,mainly,business,data,processingFragments,filled,pausesSeusss,cat,in,the,hat,is,different,from,othercats,Lemma,same,stem,part,of,speech,rough,word,sensecat,and,cats,same,lemmaWordform,the,full,inflected,surface,formcat,and,cats,different,wordformsHow,many,words,they,lay,back,on,the,San,Francisco,grass,and,looked,at,the,stars,and,theirType,an,element,of,the,vocabulary,Token,an,instance,of,that,type,in,running,text,How,many,15,tokens,or,14,13,types,or,12,or,11,How,many,words,N,number,of,tokensV,vocabulary,set,of,types,V,is,the,size,of,the,vocabularyTokens,NTypes,V,Switchboard,phoneconversations2,4,million20thousandShakespeare884,00031thousandGoogle,N,grams1,trillion13,millionChurch,and,Gale,1990,V,O,N,Simple,Tokenization,in,UNIX,Inspired,by,Ken,Churchs,UNIX,for,Poets,Given,a,text,file,output,the,word,tokens,and,their,frequenciestr,scA,Za,z,n,shakes,txt,sort,uniqc,1945,A72,AARON19,ABBESS5,ABBOT,25,Aaron6,Abate1,Abates5,Abbess6,Abbey3,Abbot,Change,all,non,alpha,tonewlinesSort,in,alphabetical,orderMerge,and,count,each,typeThe,first,step,tokenizingtr,scA,Za,z,n,shakes,txt,headTHESONNETSbyWilliamShakespeareFromfairestcreaturesWe,The,second,step,sortingtr,scA,Za,z,n,shakes,txt,sort,headAAAAAAAAA,More,countingMerging,upper,and,lower,casetrA,Z,a,z,shakes,txt,tr,scA,Za,z,n,sort,uniqc,Sorting,the,countstrA,Z,a,z,shakes,txt,tr,scA,Za,z,n,sort,uniqc,sort,n,r23243,the22225,i18618,and16339,to15687,of12780,a12163,you10839,my10005,in8954,dWhat,happened,here,Issues,in,TokenizationFinlands,capital,Finland,FinlandsFinlands,whatre,Im,isntWhat,are,I,am,is,notHewlett,Packard,Hewlett,Packard,state,of,the,art,state,of,the,art,Lowercaselower,case,lowercase,lower,case,San,Franciscoone,token,or,two,m,p,h,PhD,Tokenization,language,issuesFrenchL,ensembleone,token,or,two,L,L,Le,Want,lensembleto,match,with,un,ensembleGerman,noun,compounds,are,not,segmentedLebensversicherungsgesellschaftsangestellterlife,insurance,company,employeeGerman,information,retrieval,needs,compound,splitterTokenization,language,issuesChinese,and,Japanese,no,spaces,between,words,Sharapovanow,lives,in,US,southeastern,FloridaFurther,complicated,in,Japanese,with,multiple,alphabets,intermingledDates,amounts,in,multiple,formats500,500K,6,000,KatakanaHiraganaKanjiRomajiEnd,user,can,express,query,entirely,in,hiragana,Word,Tokenization,in,ChineseAlso,called,Word,SegmentationChinese,words,are,composed,of,charactersCharacters,are,generally,1,syllable,and,1,morpheme,Average,word,is,2,4,characters,long,Standard,baseline,segmentation,algorithm,Maximum,Matching,also,called,Greedy,Maximum,MatchingWord,Segmentation,AlgorithmGiven,a,wordlist,of,Chinese,and,a,string,1,Start,a,pointer,at,the,beginning,of,the,string2,Find,the,longest,word,in,dictionary,that,matches,the,string,starting,at,pointer3,Move,the,pointer,over,the,word,in,string4,Go,to,2Max,match,segmentation,illustrationThecatinthehatThetabledownthereDoesnt,generally,work,in,English,But,works,astonishingly,well,in,ChineseModern,probabilistic,segmentation,algorithms,even,betterthe,table,down,therethe,cat,in,the,hattheta,bled,own,thereBasic,Text,ProcessingWord,tokenizationBasic,Text,ProcessingWord,Normalization,and,StemmingNormalizationNeed,to,normalize,terms,Information,Retrieval,indexed,text,query,terms,must,have,same,form,We,want,to,match,U,S,A,and,USAWe,implicitly,define,equivalence,classes,of,termse,g,deleting,periods,in,a,termAlternative,asymmetric,expansion,Enter,windowSearch,window,windowsEnter,windowsSearch,Windows,windows,windowEnter,WindowsSearch,WindowsPotentially,more,powerful,but,less,efficientCase,foldingApplications,like,IR,reduce,all,letters,to,lower,caseSince,users,tend,to,use,lower,casePossible,exception,upper,case,in,mid,sentence,e,g,General,MotorsFedvs,fedSAILvs,sailFor,sentiment,analysis,MT,Information,extractionCase,is,helpful,USversus,us,is,important,LemmatizationReduce,inflections,or,variant,forms,to,base,formam,are,is,becar,cars,car,s,cars,carthe,boy,s,cars,are,different,colorsthe,boy,car,be,different,colorLemmatization,have,to,find,correct,dictionary,headword,formMachine,translationSpanish,quiero,I,want,quieres,you,want,same,lemma,as,quererwantMorphologyMorphemes,The,small,meaningful,units,that,make,up,wordsStems,The,core,meaning,bearing,unitsAffixes,Bits,and,pieces,that,adhere,to,stemsOften,with,grammatical,functionsStemmingReduce,terms,to,their,stems,in,information,retrievalStemmingis,crude,chopping,of,affixeslanguage,dependente,g,automate,s,automatic,automationall,reduced,to,automat,for,example,compressed,and,compression,are,both,accepted,as,equivalent,to,compress,for,examplcompress,andcompress,arboth,acceptas,equivalto,compressPorters,algorithmThe,most,common,English,stemmerStep,1assessscaresses,caressiesiponies,ponisssscaress,caresss,cats,catStep,1b,v,ingwalking,walksing,sing,v,edplastered,plasterStep,2,for,long,stems,ationalate,relationalrelateizerizedigitizer,digitizeatorateoperator,operateStep,3,for,longer,stems,al,revival,revivable,adjustable,adjustate,activate,activViewing,morphology,in,a,corpusWhy,only,strip,ingif,there,is,a,vowel,v,ingwalking,walksing,sing,36Viewing,morphology,in,a,corpusWhy,only,strip,ingif,there,is,a,vowel,v,ingwalking,walksing,sing,37tr,sc,A,Za,z,n,shakes,txt,greping,sort,uniq,c,sort,nr,tr,sc,A,Za,z,n,shakes,txt,grep,aeiou,ing,sort,uniq,c,sort,nr548,being541,nothing152,something145,coming130,morning122,having120,living117,loving116,Being102,going1312,King548,being541,nothing388,king375,bring358,thing307,ring152,something145,coming130,morning,Dealing,with,complex,morphology,is,sometimes,necessarySome,languages,requires,complex,morpheme,segmentationTurkishUygarlastiramadiklarimizdanmissinizcasina,behaving,as,if,you,are,among,those,whom,we,could,not,civilizeUygar,civilized,las,become,tir,cause,ama,not,able,dik,past,larplural,imizp1pl,danabl,mispast,siniz2pl,casinaas,if,Basic,Text,ProcessingWord,Normalization,and,StemmingBasic,Text,ProcessingSentence,Segmentation,and,Decision,TreesSentence,Segmentation,are,relatively,unambiguousPeriod,is,quite,ambiguousSentence,boundaryAbbreviations,like,Inc,or,Dr,Numbers,like,02,or,4,3Build,a,binary,classifierLooks,at,a,Decides,EndOfSentence,NotEndOfSentenceClassifiers,hand,written,rules,regular,expressions,or,machine,learningDetermining,if,a,word,is,end,of,sentence,a,Decision,Tree,More,sophisticated,decision,tree,featuresCase,of,word,with,Upper,Lower,Cap,NumberCase,of,word,after,Upper,Lower,Cap,NumberNumeric,featuresLength,of,word,with,Probability,word,with,occurs,at,end,of,s,Probability,word,after,occurs,at,beginning,of,s,Implementing,Decision,TreesA,decision,tree,is,just,an,if,then,else,statementThe,interesting,research,is,choosing,the,featuresSetting,up,the,structure,is,often,too,hard,to,do,by,handHand,building,only,possible,for,very,simple,features,domainsFor,numeric,features,its,too,hard,to,pick,each,thresholdInstead,structure,usually,learned,by,machine,learning,from,a,training,corpusDecision,Trees,and,other,classifiersWe,can,think,of,the,questions,in,a,decision,treeAs,features,that,could,be,exploited,by,any,kind,of,classifierLogistic,regressionSVMNeural,Netsetc,Basic,Text,ProcessingSentence,Segmentation,and,Decision,Trees'

print(use_regex(text))


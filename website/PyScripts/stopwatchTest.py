# print(stopwords.words('english'))
DIC_LOCATION = "website/static/wordDic/wordninja_words.txt"


import csv
import wordninja
import time
import re
from pathlib import Path
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize
from PyPDF2 import PdfFileReader
from werkzeug.utils import secure_filename
from os.path import join
from docx import Document
from pptx import Presentation
import re


START_TIME = time.time()
textExctLoc = Path("website/textExtract")

# def use_regex(input_text):
#     text = re.sub("[^\x00-\x7F]+","", input_text)
#     text = re.sub("\W+", " ", text)
#     return text

# def cleaning(unclean:str):
#     stop_words = set(stopwords.words('english'))
#     slicedwords = wordninja.split(unclean)
#     filtered_sentence = list(set([w.lower() for w in slicedwords if not w.lower() in stop_words]))
#     return filtered_sentence

def StrToCSV(extLocation: str, unclean:str, fileName: str):
    FileName = f"{fileName}.csv"
    text = re.sub("[^\x00-\x7F]+","", unclean)
    text = re.sub("\W+", " ", text)
    stop_words = set(stopwords.words('english'))
    slicedwords = wordninja.split(text)
    filtered_sentence = list(set([w.lower() for w in slicedwords if not w.lower() in stop_words]))
    with Path(join(extLocation,secure_filename(FileName))).open(mode="w", encoding='utf-8') as outputFile:
        writer = csv.writer(outputFile)
        writer.writerow(filtered_sentence)
        outputFile.close()
        return 1

def PdfToText (fileobj: str):
    pdf = PdfFileReader(fileobj)
    text = ''
    for page in pdf.pages:
        text += page.extract_text()
    return text

def DocxToText(fileobj: str):
    document = Document(fileobj)
    docText = '\n\n'.join(paragraph.text + " " for paragraph in document.paragraphs)
    return docText

def pptxToTxt(fileobj: str):
    prs = Presentation(fileobj)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text + " "
    return text

def txtReader(fileobj: str):

    with open(fileobj, 'r', encoding='UTF-8') as file:
        textBufffer = ''
        for line in file:
            textBufffer += line.rstrip('\n')
        return textBufffer

def pptOrDocToPDF():
    
    pass

# StrToCSV(textExctLoc,txtReader("website/PyScripts/2_TextProc.txt"),'textFile')

# print(txtReader("website/PyScripts/testtxt1.txt"))

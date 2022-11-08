import csv
import wordninja
import re
from pathlib import Path
from nltk.corpus import stopwords
from PyPDF2 import PdfFileReader
from werkzeug.utils import secure_filename
from os.path import join
from docx import Document
from pptx import Presentation
import re

def StrToCSV(extLocation: str, unclean:str, fileName: str):
    FileName = Path(fileName).stem + '.csv'
    text = re.sub("[^\x00-\x7F]+","", unclean)
    text = re.sub("\W+", " ", text)
    stop_words = set(stopwords.words('english'))
    slicedwords = wordninja.split(text)
    filtered_sentence = list(set([w.lower() for w in slicedwords if not w.lower() in stop_words]))
    with Path(join(extLocation,secure_filename(FileName))).open(mode="w", encoding='utf-8') as outputFile:
        writer = csv.writer(outputFile)
        writer.writerow(filtered_sentence)
        outputFile.close()
        return 1

def PdfToText (fileobj: object):
    pdf = PdfFileReader(fileobj)
    text = ''
    for page in pdf.pages:
        text += page.extract_text()
    return text

def DocxToText(fileobj: object):
    document = Document(fileobj)
    docText = '\n\n'.join(paragraph.text + " " for paragraph in document.paragraphs)
    return docText

def pptxToTxt(fileobj: object):
    prs = Presentation(fileobj)
    text = ''
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text += run.text + " "
    return text

def txtReader(fileobj: object):
    obj=fileobj.read()
    text = obj.decode("utf-8")
    return str(text)
